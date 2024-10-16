# :coding: utf-8

from core import Core
from db_conn import DBconn
from pptx import Presentation
from pptx.util import Inches
from pptx.enum.text import PP_ALIGN
from pptx.util import Pt
import markdown


class PPT( Core ):
    def write( self , prompt , scenario = True):
        if scenario:
            content = '시나리오'
        else:
            content = '시놉시스'
        search_msg  = '이 시스템은 영화 마케팅 전문가 입니다.'
        search_msg += '아래는  제작을 하려고 하는 시나리오 입니다'
        search_msg += '{body}'
        search_msg += ''
        search_msg += f'위의 {content}를 이용해서 영화를 만들기 위한 투자자들이 보고'
        search_msg += '투자를 결정 할수 있도록 프레젠 테이션 파일에 들어갈 슬라이드의 내용들은 아래와 같다'
        search_msg += '['
        search_msg += '영화의 소개 : [영화의 개요, 장르, 목표 관객층,],'
        search_msg += '시놉시스 : [간단한 줄거리 요약, 주요 테마와 메시지 ],'
        search_msg += '캐릭터 소개 : [ 주요 캐릭터와 그들의 역할 ], '
        search_msg += '시장 분석 : [현재 영화 시장 동향, 타겟 관객층 분석, 비슷한 장르의 성공 사례],'
        search_msg += '예산계획 : [총예산, 주요 비용 항목( 예:캐스팅, 촬영, 후반작업 등 )] , '
        search_msg += '수익 예측 : [ 예상 수익원( 박스 오피스, 스트리밍 등 ) ],'
        search_msg += '투자 혜택 : [ 투자자들에게 제공될 혜택, ROI 예측 ],'
        search_msg += '제작 일정 : [주요 마일스톤, 예상 완료 날짜],'
        search_msg += ']'
        search_msg += '이 내용들중 다른 키워드는 작성 하지 말고 아래의 내용 작성 해줘'
        search_msg += prompt
        chain = self.chain( search_msg )
        if scenario:
            response = chain.invoke({'body':self.scenario}  )
        else:
            response = chain.invoke({'body':self.synop}  )
        return response

    def raw_write( self, prompt ):
        search_msg  = '이 시스템은 영화 마케팅 전문가 입니다.'
        search_msg += prompt
        search_msg += '다른 내용은 제외 하고 결과만 출력해주세요' 
        chain = self.chain( search_msg )
        response = chain.invoke({'content':self.scenario}  )
        return response
        
def add_full_slide_text_with_font_size(presentation, title_text, content_text, font_size=18):
    # 슬라이드 추가
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])  # 빈 레이아웃

    # 슬라이드 제목 추가
    title = slide.shapes.title
    title.text = title_text
    
    # 슬라이드 크기 계산
    slide_width = presentation.slide_width
    slide_height = presentation.slide_height

    # 텍스트 상자 생성
    left = Inches(0.5)
    top = Inches(1)
    width = slide_width - Inches(1)
    height = slide_height - Inches(2)
    
    text_box = slide.shapes.add_textbox(left, top, width, height)
    text_frame = text_box.text_frame
    text_frame.text = content_text

    # 텍스트를 중앙 정렬하고 폰트 크기 설정
    for paragraph in text_frame.paragraphs:
        paragraph.alignment = PP_ALIGN.CENTER
        for run in paragraph.runs:
            run.font.size = Pt(font_size)


def write_slide( prst,slide_num, title, content ):
    slide = prst.slides[0]
    print( '\n' *3 )
    for slide in prst.slides :
        for shape in slide.shapes:
            print( shape.name )
            text_frame = shape.text_frame
            if 'title' in shape.name:
                text_frame.text = title
            elif 'textbox' in shape.name:
                text_frame.clear()
                text_frame.text = content
                
        
def main():
    ppt = PPT()
    with open( './tmp/uploaded/시나리오.txt' ) as f:
        ppt.scenario = f.read()

    ## 소개
    prompt =  '다른 내용은 포함 하지 말고 [영화의 소개]에 들어갈 내용만을 작성해주세요.' 
    prompt += '프레젠 테이션에 슬라이드에 사용할 텍스트로만  작성해줘.'
    overview = ppt.write( prompt )


    prst = Presentation( './tmp/template.pptx' )
    write_slide( prst, 0, '영화의 소개', overview )
#    slide_1 = presentation.slides.add_slide(presentation.slide_layouts[0])
#    title_1 = slide_1.shapes.title
#    content_1 = slide_1.placeholders[1]
#    title_1.text = "영화 개요"

    #overview = markdown.markdown( overview )
    #lines = overview.splitlines()

#    left = Inches( 0.5)
#    top = Inches( 1)
#    width  = presentation.slide_width - Inches( 1 )
#    height = presentation.slide_height - Inches( 2 )


    #content_1.text = overview 
    #add_full_slide_text_with_font_size(presentation, "영화 개요", overview, font_size=14)
    #pprint( overview )        
    #presentation.save( './tmp/test_ppt.pptx' )
    return

    ## 시놉시스
    ppt.synop = DBconn().last_synop()
    synop_theme = ppt.write( '시놉시스를 보고 주요 테마와 메시지를 작성해 주세요.', False )

    ## 캐릭터 소개
    char    = ppt.write( '주요 캐릭터와 그들의 역할에 대해 작성해주세요.' )

    ## 시장 분석
    market  = ppt.raw_write( '현재 영화 시장 동향에 대해 분석해 주세요.' )
    target  = ppt.write( '타겟 관객층에 대해 분석해 주세요.' )
    genre   = ppt.write( '비슷한 장르의 성공사례에 대해 분석해 주세요.' )

    ## 예산 계획
    total = ppt.write( '총 예산을 분석해 주세요.' )

    ## 수익 예측
    predict_value = ppt.write( '극장 수익, 스트리밍 등을 통한 예상 수익을 분석해 주세요' )
    
    ## 투자 혜택
    invest = ppt.write( '투자자들에게 제공될 혜택에 대해 예측하고 분석해 주세요.' )
    roi = ppt.write( 'ROI에 대해 예측 하고 분석 해 주세요.' )

    ## 제작일정
    milestone = ppt.write( '주요 마일스톤에대해 분석하고 정리해주세요.' )
    duedate = ppt.write( '예상 완료 날짜에 대해 예측하고 분석해 주세요.' )



if __name__ == '__main__':
    from pprint import pprint
    main()

